import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import docx2txt
import fitz  # PyMuPDF

def extract_pptx(filepath):
    prs = Presentation(filepath)
    text_runs = []
    
    print(f"\n--- Extracting PPTX: {filepath} ---")
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i + 1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                print(shape.text.strip())
            
            # Identify tables
            if shape.has_table:
                print("[TABLE FOUND]")
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    print(" | ".join(row_data))
                    
            # Identify charts
            if shape.has_chart:
                print("[CHART FOUND]")
                chart = shape.chart
                print(f"Chart Type: {chart.chart_type}")
                if chart.has_title:
                    print(f"Chart Title: {chart.chart_title.text_frame.text}")

def extract_docx(filepath):
    print(f"\n--- Extracting DOCX: {filepath} ---")
    text = docx2txt.process(filepath)
    print(text)

def extract_pdf(filepath):
    print(f"\n--- Extracting PDF: {filepath} ---")
    doc = fitz.open(filepath)
    for i, page in enumerate(doc):
        print(f"\nPage {i + 1}:")
        print(page.get_text())

if __name__ == "__main__":
    for arg in sys.argv[1:]:
        if os.path.isdir(arg):
            for file in os.listdir(arg):
                path = os.path.join(arg, file)
                if file.endswith(".pptx") and not file.startswith("~"):
                    extract_pptx(path)
                elif file.endswith(".docx") and not file.startswith("~"):
                    extract_docx(path)
                elif file.endswith(".pdf") and not file.startswith("~"):
                    extract_pdf(path)
        else:
            if arg.endswith(".pptx") and not arg.startswith("~"):
                extract_pptx(arg)
            elif arg.endswith(".docx") and not arg.startswith("~"):
                extract_docx(arg)
            elif arg.endswith(".pdf") and not arg.startswith("~"):
                extract_pdf(arg)
